"""
Manava Analytics Server
Run: python3 server.py
Listens on http://127.0.0.1:5050
"""

from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import pandas as pd
import numpy as np
import json, io, os, traceback, warnings
from datetime import datetime

warnings.filterwarnings("ignore")

app = Flask(__name__)
CORS(app, origins="*", allow_headers="*", methods=["GET","POST","OPTIONS"])

# ── Helpers ────────────────────────────────────────────────────────────────────

def parse_data(payload):
    """Parse incoming data from various sources into a DataFrame."""
    src = payload.get("source")
    raw = payload.get("data")

    if src == "csv_text":
        return pd.read_csv(io.StringIO(raw))
    elif src == "json":
        return pd.DataFrame(json.loads(raw))
    elif src == "tsv":
        return pd.read_csv(io.StringIO(raw), sep="\t")
    else:
        # Try CSV first, then TSV
        try:
            return pd.read_csv(io.StringIO(raw))
        except Exception:
            return pd.read_csv(io.StringIO(raw), sep="\t")


def safe_val(v):
    """Convert numpy types to Python native for JSON serialisation."""
    if isinstance(v, (np.integer,)):  return int(v)
    if isinstance(v, (np.floating,)): return None if np.isnan(v) else float(v)
    if isinstance(v, (np.ndarray,)):  return v.tolist()
    if isinstance(v, pd.Timestamp):   return str(v)
    if pd.isna(v):                    return None
    return v


def df_info(df):
    """Return basic dataset metadata."""
    return {
        "rows": len(df),
        "columns": len(df.columns),
        "column_names": df.columns.tolist(),
        "dtypes": {c: str(df[c].dtype) for c in df.columns},
        "missing": {c: int(df[c].isna().sum()) for c in df.columns},
        "sample": df.head(5).map(safe_val).to_dict(orient="records"),
    }


# ── Analysis Engines ───────────────────────────────────────────────────────────

def descriptive_analysis(df):
    num = df.select_dtypes(include="number")
    cat = df.select_dtypes(exclude="number")

    result = {"type": "descriptive", "numeric": {}, "categorical": {}, "overview": df_info(df)}

    for col in num.columns:
        s = num[col].dropna()
        result["numeric"][col] = {
            "mean":   safe_val(s.mean()),
            "median": safe_val(s.median()),
            "std":    safe_val(s.std()),
            "min":    safe_val(s.min()),
            "max":    safe_val(s.max()),
            "q25":    safe_val(s.quantile(0.25)),
            "q75":    safe_val(s.quantile(0.75)),
            "skew":   safe_val(s.skew()),
            "count":  int(s.count()),
            "values": s.tolist()[:200],   # for charts
        }

    for col in cat.columns:
        vc = df[col].value_counts().head(20)
        result["categorical"][col] = {
            "unique":  int(df[col].nunique()),
            "top":     str(df[col].mode().iloc[0]) if len(df[col].dropna()) else None,
            "counts":  {str(k): int(v) for k, v in vc.items()},
        }

    return result


def correlation_analysis(df):
    num = df.select_dtypes(include="number")
    if num.shape[1] < 2:
        return {"type": "correlation", "error": "Need at least 2 numeric columns"}

    corr = num.corr()
    pairs = []
    cols = corr.columns.tolist()
    for i in range(len(cols)):
        for j in range(i+1, len(cols)):
            v = safe_val(corr.iloc[i, j])
            if v is not None:
                pairs.append({"col_a": cols[i], "col_b": cols[j], "r": round(v, 4)})

    pairs.sort(key=lambda x: abs(x["r"]), reverse=True)

    return {
        "type": "correlation",
        "matrix": {c: {r: safe_val(corr.loc[c, r]) for r in cols} for c in cols},
        "columns": cols,
        "top_pairs": pairs[:10],
        "scatter_data": {
            col: num[col].dropna().tolist()[:200] for col in cols
        },
    }


def forecasting_analysis(df):
    num = df.select_dtypes(include="number")
    results = {"type": "forecasting", "forecasts": {}}

    for col in num.columns[:3]:
        s = num[col].dropna().reset_index(drop=True)
        if len(s) < 4:
            continue
        try:
            n_pred = max(5, len(s) // 5)
            # Exponential smoothing (manual, no scipy dependency)
            alpha = 0.3
            smoothed = [float(s.iloc[0])]
            for i in range(1, len(s)):
                smoothed.append(alpha * float(s.iloc[i]) + (1 - alpha) * smoothed[-1])
            # Trend from last few points
            trend = (smoothed[-1] - smoothed[-min(5,len(smoothed))]) / min(5,len(smoothed))
            forecast = [smoothed[-1] + trend * (i + 1) for i in range(n_pred)]
            results["forecasts"][col] = {
                "historical":   [safe_val(v) for v in s.tolist()],
                "forecast":     [safe_val(v) for v in forecast],
                "n_historical": len(s),
                "n_forecast":   n_pred,
                "method":       "exponential_smoothing",
            }
        except Exception as e:
            x = np.arange(len(s))
            p = np.polyfit(x, s, 1)
            xf = np.arange(len(s), len(s) + 5)
            results["forecasts"][col] = {
                "historical":   [safe_val(v) for v in s.tolist()],
                "forecast":     [safe_val(v) for v in np.polyval(p, xf)],
                "n_historical": len(s),
                "n_forecast":   5,
                "method":       "linear",
            }

    return results


def clustering_analysis(df):
    from sklearn.preprocessing import StandardScaler
    from sklearn.cluster import KMeans
    from sklearn.decomposition import PCA

    num = df.select_dtypes(include="number").dropna()
    if num.shape[1] < 1 or num.shape[0] < 4:
        return {"type": "clustering", "error": "Need at least 4 rows and 1 numeric column"}

    scaler = StandardScaler()
    X = scaler.fit_transform(num)

    # Elbow method — try k=2..6
    inertias = []
    k_range  = range(2, min(7, len(num)))
    for k in k_range:
        km = KMeans(n_clusters=k, random_state=42, n_init=10)
        km.fit(X)
        inertias.append(safe_val(km.inertia_))

    best_k = 3 if len(num) > 10 else 2
    km = KMeans(n_clusters=best_k, random_state=42, n_init=10)
    labels = km.fit_predict(X)

    # PCA for 2-D scatter
    pca = PCA(n_components=min(2, X.shape[1]))
    X2d = pca.fit_transform(X)

    points = []
    for i in range(len(X2d)):
        row = {"cluster": int(labels[i]), "x": safe_val(X2d[i, 0])}
        row["y"] = safe_val(X2d[i, 1]) if X2d.shape[1] > 1 else 0
        # attach original values
        for c in num.columns:
            row[c] = safe_val(num.iloc[i][c])
        points.append(row)

    # Cluster summaries
    num2 = num.copy()
    num2["_cluster"] = labels
    summaries = {}
    for k in range(best_k):
        grp = num2[num2["_cluster"] == k].drop("_cluster", axis=1)
        summaries[k] = {c: {"mean": safe_val(grp[c].mean()), "count": int(len(grp))} for c in grp.columns}

    return {
        "type":      "clustering",
        "k":         best_k,
        "points":    points,
        "summaries": summaries,
        "inertias":  inertias,
        "k_range":   list(k_range),
        "explained_variance": [safe_val(v) for v in pca.explained_variance_ratio_],
    }


# ── Routes ─────────────────────────────────────────────────────────────────────

@app.route("/ping")
def ping():
    return jsonify({"status": "ok", "server": "Manava Analytics"})


@app.route("/analyze", methods=["POST"])
def analyze():
    try:
        body   = request.json
        df     = parse_data(body)
        mode   = body.get("mode", "all")
        result = {"columns": df.columns.tolist(), "rows": len(df), "analyses": {}}

        if mode in ("all", "descriptive"):
            result["analyses"]["descriptive"] = descriptive_analysis(df)
        if mode in ("all", "correlation"):
            result["analyses"]["correlation"] = correlation_analysis(df)
        if mode in ("all", "forecasting"):
            result["analyses"]["forecasting"] = forecasting_analysis(df)
        if mode in ("all", "clustering"):
            result["analyses"]["clustering"] = clustering_analysis(df)

        return jsonify(result)

    except Exception:
        return jsonify({"error": traceback.format_exc()}), 500


@app.route("/export/excel", methods=["POST"])
def export_excel():
    try:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, LineChart, Reference
        from openpyxl.chart.series import SeriesLabel
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        body = request.json
        df   = parse_data(body)
        data = body.get("analysis", {})

        wb = Workbook()

        # ── Colour palette ─────────────────────────────────────────────────
        PURPLE = "6C63FF"
        DARK   = "13112D"
        GOLD   = "C9A84C"
        WHITE  = "FFFFFF"
        GRAY   = "F4F3FF"

        header_font  = Font(name="Calibri", bold=True, color=WHITE, size=12)
        header_fill  = PatternFill("solid", fgColor=PURPLE)
        subhdr_fill  = PatternFill("solid", fgColor="9D97FF")
        subhdr_font  = Font(name="Calibri", bold=True, color=WHITE, size=11)
        title_font   = Font(name="Calibri", bold=True, color=DARK, size=14)
        normal_font  = Font(name="Calibri", size=11)
        center_align = Alignment(horizontal="center", vertical="center")
        thin_border  = Border(
            left=Side(style="thin", color="DDDDFF"),
            right=Side(style="thin", color="DDDDFF"),
            top=Side(style="thin", color="DDDDFF"),
            bottom=Side(style="thin", color="DDDDFF"),
        )

        def style_header(ws, row, cols):
            for c in range(1, cols + 1):
                cell = ws.cell(row=row, column=c)
                cell.font  = header_font
                cell.fill  = header_fill
                cell.alignment = center_align
                cell.border = thin_border

        def style_row(ws, row, cols, alt=False):
            fill = PatternFill("solid", fgColor="EEEDFE" if alt else WHITE)
            for c in range(1, cols + 1):
                cell = ws.cell(row=row, column=c)
                cell.font   = normal_font
                cell.fill   = fill
                cell.border = thin_border
                cell.alignment = Alignment(vertical="center")

        # ── Sheet 1: Raw Data ──────────────────────────────────────────────
        ws1 = wb.active
        ws1.title = "Raw Data"
        ws1.sheet_view.showGridLines = False
        ws1.row_dimensions[1].height = 28

        ws1.cell(1, 1, "Manava Analytics — Raw Data").font = title_font
        ws1.cell(1, 1).alignment = center_align
        ws1.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(df.columns))

        for ci, col in enumerate(df.columns, 1):
            cell = ws1.cell(2, ci, col)
            cell.font = header_font; cell.fill = header_fill
            cell.alignment = center_align; cell.border = thin_border
            ws1.column_dimensions[get_column_letter(ci)].width = max(14, len(str(col)) + 4)

        for ri, row_data in enumerate(df.itertuples(index=False), 3):
            for ci, val in enumerate(row_data, 1):
                ws1.cell(ri, ci, val).border = thin_border
                ws1.cell(ri, ci).font = normal_font
                if ri % 2 == 0:
                    ws1.cell(ri, ci).fill = PatternFill("solid", fgColor="EEEDFE")

        # ── Sheet 2: Descriptive Stats ─────────────────────────────────────
        ws2 = wb.create_sheet("Descriptive Stats")
        ws2.sheet_view.showGridLines = False
        desc = data.get("analyses", {}).get("descriptive", {})
        num_data = desc.get("numeric", {})

        ws2.row_dimensions[1].height = 28
        ws2.cell(1, 1, "Descriptive Statistics").font = title_font
        ws2.cell(1, 1).alignment = center_align
        ws2.merge_cells(start_row=1, start_column=1, end_row=1, end_column=9)

        headers = ["Column", "Count", "Mean", "Median", "Std Dev", "Min", "Max", "Q25", "Q75"]
        for ci, h in enumerate(headers, 1):
            ws2.cell(2, ci, h)
        style_header(ws2, 2, len(headers))

        for ri, (col, stats) in enumerate(num_data.items(), 3):
            vals = [col, stats.get("count"), round(stats.get("mean") or 0, 4),
                    round(stats.get("median") or 0, 4), round(stats.get("std") or 0, 4),
                    round(stats.get("min") or 0, 4), round(stats.get("max") or 0, 4),
                    round(stats.get("q25") or 0, 4), round(stats.get("q75") or 0, 4)]
            for ci, v in enumerate(vals, 1):
                ws2.cell(ri, ci, v)
            style_row(ws2, ri, len(headers), alt=(ri % 2 == 0))

        for i in range(1, 10):
            ws2.column_dimensions[get_column_letter(i)].width = 14

        # Add bar chart for means
        if num_data:
            chart_data_row = 2
            chart_data_col = 1
            chart = BarChart()
            chart.type  = "col"
            chart.title = "Column Means"
            chart.style = 10
            chart.y_axis.title = "Mean Value"
            chart.x_axis.title = "Column"

            n_cols = len(num_data)
            data_ref = Reference(ws2, min_col=3, min_row=2, max_row=2 + n_cols)
            cats_ref = Reference(ws2, min_col=1, min_row=3, max_row=2 + n_cols)
            chart.add_data(data_ref, titles_from_data=True)
            chart.set_categories(cats_ref)
            chart.height = 14; chart.width = 22
            ws2.add_chart(chart, f"K3")

        # ── Sheet 3: Correlation ───────────────────────────────────────────
        ws3 = wb.create_sheet("Correlation")
        ws3.sheet_view.showGridLines = False
        corr_data = data.get("analyses", {}).get("correlation", {})
        cols_list = corr_data.get("columns", [])
        matrix    = corr_data.get("matrix", {})

        ws3.cell(1, 1, "Correlation Matrix").font = title_font
        ws3.cell(1, 1).alignment = center_align
        if cols_list:
            ws3.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(cols_list)+1)

        for ci, col in enumerate(cols_list, 2):
            ws3.cell(2, ci, col)
        ws3.cell(2, 1, "")
        style_header(ws3, 2, len(cols_list) + 1)

        for ri, rc in enumerate(cols_list, 3):
            ws3.cell(ri, 1, rc).font = header_font
            ws3.cell(ri, 1).fill = subhdr_fill
            ws3.cell(ri, 1).alignment = center_align
            for ci, cc in enumerate(cols_list, 2):
                v = matrix.get(rc, {}).get(cc)
                cell = ws3.cell(ri, ci, round(v, 3) if v is not None else "")
                cell.font = normal_font
                cell.alignment = center_align
                cell.border = thin_border
                if v is not None:
                    abs_v = abs(v)
                    if abs_v > 0.7:
                        cell.fill = PatternFill("solid", fgColor="AFA9EC")
                    elif abs_v > 0.4:
                        cell.fill = PatternFill("solid", fgColor="CEEBFC")
                    else:
                        cell.fill = PatternFill("solid", fgColor=WHITE)

        for i in range(1, len(cols_list) + 2):
            ws3.column_dimensions[get_column_letter(i)].width = 14

        # ── Sheet 4: Forecast ──────────────────────────────────────────────
        ws4 = wb.create_sheet("Forecast")
        ws4.sheet_view.showGridLines = False
        fc_data = data.get("analyses", {}).get("forecasting", {}).get("forecasts", {})

        ws4.cell(1, 1, "Forecast Analysis").font = title_font
        ws4.merge_cells(start_row=1, start_column=1, end_row=1, end_column=6)

        col_offset = 1
        for col_name, fc in fc_data.items():
            hist = fc.get("historical", [])
            pred = fc.get("forecast", [])

            ws4.cell(2, col_offset, col_name).font = subhdr_font
            ws4.cell(2, col_offset).fill = subhdr_fill
            ws4.cell(2, col_offset + 1, "Type").font = subhdr_font
            ws4.cell(2, col_offset + 1).fill = subhdr_fill

            for ri, v in enumerate(hist, 3):
                ws4.cell(ri, col_offset, safe_val(v))
                ws4.cell(ri, col_offset + 1, "Historical")
            for ri, v in enumerate(pred, 3 + len(hist)):
                ws4.cell(ri, col_offset, safe_val(v))
                ws4.cell(ri, col_offset + 1, "Forecast")
                ws4.cell(ri, col_offset).fill = PatternFill("solid", fgColor="CEEBFC")

            # Line chart
            lc = LineChart()
            lc.title  = f"Forecast — {col_name}"
            lc.style  = 10
            lc.y_axis.title = col_name
            lc.x_axis.title = "Period"
            total_rows = len(hist) + len(pred)
            data_ref = Reference(ws4, min_col=col_offset, min_row=2, max_row=2 + total_rows)
            lc.add_data(data_ref, titles_from_data=True)
            lc.height = 14; lc.width = 22
            ws4.add_chart(lc, f"{get_column_letter(col_offset + 3)}3")
            col_offset += 14

        # ── Sheet 5: Clustering ────────────────────────────────────────────
        ws5 = wb.create_sheet("Clustering")
        ws5.sheet_view.showGridLines = False
        cl_data  = data.get("analyses", {}).get("clustering", {})
        summaries= cl_data.get("summaries", {})
        k        = cl_data.get("k", 0)

        ws5.cell(1, 1, f"Cluster Analysis  (k={k})").font = title_font
        ws5.merge_cells(start_row=1, start_column=1, end_row=1, end_column=6)

        row_idx = 3
        for cluster_id, cols_info in summaries.items():
            ws5.cell(row_idx, 1, f"Cluster {int(cluster_id) + 1}").font = subhdr_font
            ws5.cell(row_idx, 1).fill = subhdr_fill
            ws5.cell(row_idx, 2, "Mean").font = subhdr_font
            ws5.cell(row_idx, 2).fill = subhdr_fill
            ws5.cell(row_idx, 3, "Count").font = subhdr_font
            ws5.cell(row_idx, 3).fill = subhdr_fill
            row_idx += 1
            for col_name, stats in cols_info.items():
                ws5.cell(row_idx, 1, col_name).font = normal_font
                ws5.cell(row_idx, 2, round(stats.get("mean") or 0, 4)).font = normal_font
                ws5.cell(row_idx, 3, stats.get("count")).font = normal_font
                row_idx += 1
            row_idx += 1

        for i in range(1, 7):
            ws5.column_dimensions[get_column_letter(i)].width = 16

        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)
        return send_file(buf, mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                         as_attachment=True, download_name=f"manava_report_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx")

    except Exception:
        return jsonify({"error": traceback.format_exc()}), 500


@app.route("/export/pptx", methods=["POST"])
def export_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE

        body = request.json
        df   = parse_data(body)
        data = body.get("analysis", {})

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        PURPLE = RGBColor(0x6C, 0x63, 0xFF)
        DARK   = RGBColor(0x13, 0x11, 0x2D)
        WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT  = RGBColor(0xE8, 0xE7, 0xFF)
        GOLD   = RGBColor(0xC9, 0xA8, 0x4C)

        blank = prs.slide_layouts[6]

        def add_slide():
            return prs.slides.add_slide(blank)

        def bg(slide, color=RGBColor(0x07, 0x05, 0x1A)):
            fill = slide.background.fill
            fill.solid(); fill.fore_color.rgb = color

        def add_textbox(slide, text, l, t, w, h, font_size=18, bold=False,
                        color=WHITE, align=PP_ALIGN.LEFT, italic=False):
            txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf  = txb.text_frame; tf.word_wrap = True
            p   = tf.paragraphs[0]; p.alignment = align
            run = p.add_run(); run.text = text
            run.font.size   = Pt(font_size)
            run.font.bold   = bold
            run.font.italic = italic
            run.font.color.rgb = color
            return txb

        def gold_bar(slide, t=0.9):
            bar = slide.shapes.add_shape(1, Inches(0.5), Inches(t), Inches(12.33), Inches(0.04))
            bar.fill.solid(); bar.fill.fore_color.rgb = GOLD
            bar.line.fill.background()

        # ── Slide 1: Title ─────────────────────────────────────────────────
        s1 = add_slide(); bg(s1)
        gold_bar(s1, 0.8)
        add_textbox(s1, "MANAVA", 0.5, 0.2, 12, 0.7, 48, True, PURPLE, PP_ALIGN.CENTER)
        add_textbox(s1, "Data Analytics Report", 0.5, 1.0, 12, 0.7, 28, False, LIGHT, PP_ALIGN.CENTER, italic=True)
        add_textbox(s1, f"Generated {datetime.now().strftime('%B %d, %Y')}", 0.5, 1.8, 12, 0.4, 14, False, RGBColor(0x9D, 0x97, 0xFF), PP_ALIGN.CENTER)

        desc_txt = (f"Dataset: {len(df)} rows × {len(df.columns)} columns\n"
                    f"Columns: {', '.join(df.columns.tolist()[:8])}")
        add_textbox(s1, desc_txt, 1.5, 2.8, 10, 1.2, 14, False, LIGHT, PP_ALIGN.CENTER)

        # ── Slide 2: Dataset Overview ──────────────────────────────────────
        s2 = add_slide(); bg(s2)
        gold_bar(s2)
        add_textbox(s2, "Dataset Overview", 0.5, 0.2, 12, 0.6, 28, True, PURPLE, PP_ALIGN.LEFT)

        overview_lines = [
            f"Total Rows: {len(df)}",
            f"Total Columns: {len(df.columns)}",
            f"Numeric Columns: {len(df.select_dtypes(include='number').columns)}",
            f"Text Columns: {len(df.select_dtypes(exclude='number').columns)}",
            f"Missing Values: {int(df.isna().sum().sum())}",
        ]
        add_textbox(s2, "\n".join(overview_lines), 0.5, 1.2, 6, 3, 16, False, LIGHT)

        cols_detail = "\n".join([f"• {c}  [{df[c].dtype}]" for c in df.columns[:12]])
        add_textbox(s2, "Column Details:\n" + cols_detail, 7, 1.2, 6, 4, 13, False, RGBColor(0xC8, 0xC5, 0xFF))

        # ── Slide 3: Descriptive Stats ─────────────────────────────────────
        s3 = add_slide(); bg(s3)
        gold_bar(s3)
        add_textbox(s3, "Descriptive Statistics", 0.5, 0.2, 12, 0.6, 28, True, PURPLE)

        desc = data.get("analyses", {}).get("descriptive", {})
        num_cols = list(desc.get("numeric", {}).keys())

        if num_cols:
            # Chart: means
            chart_data = ChartData()
            chart_data.categories = num_cols[:8]
            chart_data.add_series("Mean", [round(desc["numeric"][c].get("mean") or 0, 2) for c in num_cols[:8]])
            chart = s3.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(0.5), Inches(1.1), Inches(7), Inches(5.5), chart_data).chart
            chart.has_title = True; chart.chart_title.text_frame.text = "Column Means"
            chart.plots[0].series[0].format.fill.solid()
            chart.plots[0].series[0].format.fill.fore_color.rgb = PURPLE

            stats_txt = "\n".join([
                f"{c}: mean={round(desc['numeric'][c].get('mean') or 0, 2)}  "
                f"std={round(desc['numeric'][c].get('std') or 0, 2)}"
                for c in num_cols[:6]
            ])
            add_textbox(s3, stats_txt, 8, 1.2, 5, 5, 13, False, LIGHT)

        # ── Slide 4: Correlation ───────────────────────────────────────────
        s4 = add_slide(); bg(s4)
        gold_bar(s4)
        add_textbox(s4, "Correlation Analysis", 0.5, 0.2, 12, 0.6, 28, True, PURPLE)

        corr = data.get("analyses", {}).get("correlation", {})
        pairs = corr.get("top_pairs", [])
        if pairs:
            pairs_txt = "Top Correlations:\n\n" + "\n".join([
                f"  {p['col_a']}  ↔  {p['col_b']}    r = {p['r']}"
                for p in pairs[:8]
            ])
            add_textbox(s4, pairs_txt, 0.5, 1.2, 7, 5, 15, False, LIGHT)

            strong = [p for p in pairs if abs(p["r"]) > 0.5]
            insight = f"Found {len(strong)} strong correlation(s) (|r| > 0.5)" if strong else "No strong linear correlations detected."
            add_textbox(s4, "Insight: " + insight, 8, 1.2, 5, 1.5, 14, False, RGBColor(0xC8, 0xC5, 0xFF), italic=True)

        # ── Slide 5: Forecasting ───────────────────────────────────────────
        s5 = add_slide(); bg(s5)
        gold_bar(s5)
        add_textbox(s5, "Forecast Analysis", 0.5, 0.2, 12, 0.6, 28, True, PURPLE)

        fc = data.get("analyses", {}).get("forecasting", {}).get("forecasts", {})
        if fc:
            first_col = list(fc.keys())[0]
            fc_info = fc[first_col]
            hist = fc_info.get("historical", [])
            pred = fc_info.get("forecast", [])

            chart_data = ChartData()
            chart_data.categories = [f"P{i+1}" for i in range(len(hist) + len(pred))]
            chart_data.add_series("Historical", hist + [None]*len(pred))
            chart_data.add_series("Forecast",   [None]*len(hist) + pred)
            chart = s5.shapes.add_chart(XL_CHART_TYPE.LINE,
                Inches(0.5), Inches(1.1), Inches(9), Inches(5.5), chart_data).chart
            chart.has_title = True; chart.chart_title.text_frame.text = f"Forecast — {first_col}"

            fc_txt = (f"Column: {first_col}\n"
                      f"Historical points: {len(hist)}\n"
                      f"Forecast points: {len(pred)}\n\n"
                      f"Next values:\n" +
                      "\n".join([f"  {round(v,2) if v else 'N/A'}" for v in pred[:5]]))
            add_textbox(s5, fc_txt, 10, 1.2, 3, 5, 13, False, LIGHT)

        # ── Slide 6: Clustering ────────────────────────────────────────────
        s6 = add_slide(); bg(s6)
        gold_bar(s6)
        cl = data.get("analyses", {}).get("clustering", {})
        add_textbox(s6, f"Cluster Analysis  (k = {cl.get('k', '?')})", 0.5, 0.2, 12, 0.6, 28, True, PURPLE)

        summaries = cl.get("summaries", {})
        if summaries:
            cl_txt_parts = []
            for cid, cols_info in summaries.items():
                lines = [f"Cluster {int(cid)+1}"]
                for col_name, stats in list(cols_info.items())[:4]:
                    lines.append(f"  {col_name}: mean = {round(stats.get('mean') or 0, 2)}  (n={stats.get('count')})")
                cl_txt_parts.append("\n".join(lines))
            add_textbox(s6, "\n\n".join(cl_txt_parts), 0.5, 1.2, 8, 5.5, 14, False, LIGHT)

            ev = cl.get("explained_variance", [])
            ev_txt = f"PCA explained variance: {round(ev[0]*100,1) if ev else '?'}%"
            add_textbox(s6, ev_txt, 9.5, 1.2, 3.5, 1, 13, False, RGBColor(0xC8, 0xC5, 0xFF), italic=True)

        buf = io.BytesIO()
        prs.save(buf); buf.seek(0)
        return send_file(buf, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                         as_attachment=True, download_name=f"manava_report_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx")

    except Exception:
        return jsonify({"error": traceback.format_exc()}), 500


@app.route("/export/pdf", methods=["POST"])
def export_pdf():
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                         Table, TableStyle, HRFlowable)
        from reportlab.lib.enums import TA_CENTER, TA_LEFT

        body = request.json
        df   = parse_data(body)
        data = body.get("analysis", {})
        buf  = io.BytesIO()

        doc = SimpleDocTemplate(buf, pagesize=A4,
                                leftMargin=2*cm, rightMargin=2*cm,
                                topMargin=2*cm, bottomMargin=2*cm)

        PURPLE_RL = colors.HexColor("#6C63FF")
        DARK_RL   = colors.HexColor("#13112D")
        LIGHT_RL  = colors.HexColor("#E8E7FF")
        GOLD_RL   = colors.HexColor("#C9A84C")

        styles = getSampleStyleSheet()
        title_style  = ParagraphStyle("Title",  parent=styles["Title"],
                                      textColor=PURPLE_RL, fontSize=26, spaceAfter=4)
        h1_style     = ParagraphStyle("H1",     parent=styles["Heading1"],
                                      textColor=PURPLE_RL, fontSize=16, spaceAfter=6)
        h2_style     = ParagraphStyle("H2",     parent=styles["Heading2"],
                                      textColor=DARK_RL,   fontSize=13, spaceAfter=4)
        body_style   = ParagraphStyle("Body",   parent=styles["Normal"],
                                      fontSize=10, leading=16, spaceAfter=4)
        caption_style= ParagraphStyle("Cap",    parent=styles["Normal"],
                                      fontSize=9, textColor=colors.grey, spaceAfter=2)

        def hr():
            return HRFlowable(width="100%", thickness=1, color=GOLD_RL, spaceAfter=8)

        def tbl(data_rows, col_widths=None):
            t = Table(data_rows, colWidths=col_widths, repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND",  (0,0), (-1,0),  PURPLE_RL),
                ("TEXTCOLOR",   (0,0), (-1,0),  colors.white),
                ("FONTNAME",    (0,0), (-1,0),  "Helvetica-Bold"),
                ("FONTSIZE",    (0,0), (-1,0),  10),
                ("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white, LIGHT_RL]),
                ("FONTSIZE",    (0,1), (-1,-1), 9),
                ("GRID",        (0,0), (-1,-1), 0.4, colors.HexColor("#DDDDFF")),
                ("VALIGN",      (0,0), (-1,-1), "MIDDLE"),
                ("TOPPADDING",  (0,0), (-1,-1), 5),
                ("BOTTOMPADDING",(0,0),(-1,-1), 5),
            ]))
            return t

        story = []

        # Title page content
        story += [
            Paragraph("MANAVA", title_style),
            Paragraph("Data Analytics Report", ParagraphStyle("sub", parent=styles["Normal"],
                fontSize=14, textColor=DARK_RL, spaceAfter=2)),
            Paragraph(datetime.now().strftime("%B %d, %Y"), caption_style),
            hr(), Spacer(1, 0.3*cm),
        ]

        # Overview
        story += [
            Paragraph("Dataset Overview", h1_style),
            Paragraph(f"Rows: {len(df)}  |  Columns: {len(df.columns)}  |  "
                      f"Missing values: {int(df.isna().sum().sum())}", body_style),
            Spacer(1, 0.3*cm),
        ]

        overview_rows = [["Column", "Type", "Missing"]]
        for col in df.columns:
            overview_rows.append([col, str(df[col].dtype), str(int(df[col].isna().sum()))])
        story += [tbl(overview_rows, [8*cm, 5*cm, 3*cm]), Spacer(1, 0.5*cm)]

        # Descriptive
        story += [hr(), Paragraph("Descriptive Statistics", h1_style)]
        desc = data.get("analyses", {}).get("descriptive", {})
        num_data = desc.get("numeric", {})
        if num_data:
            rows = [["Column", "Count", "Mean", "Median", "Std Dev", "Min", "Max"]]
            for col, s in num_data.items():
                rows.append([col,
                    str(s.get("count","")),
                    f"{s.get('mean') or 0:.4f}",
                    f"{s.get('median') or 0:.4f}",
                    f"{s.get('std') or 0:.4f}",
                    f"{s.get('min') or 0:.4f}",
                    f"{s.get('max') or 0:.4f}"])
            story += [tbl(rows), Spacer(1, 0.5*cm)]

        # Correlation
        corr = data.get("analyses", {}).get("correlation", {})
        pairs = corr.get("top_pairs", [])
        if pairs:
            story += [hr(), Paragraph("Correlation Analysis", h1_style)]
            rows = [["Column A", "Column B", "Pearson r", "Strength"]]
            for p in pairs[:10]:
                r = p["r"]
                strength = "Strong" if abs(r)>0.7 else "Moderate" if abs(r)>0.4 else "Weak"
                rows.append([p["col_a"], p["col_b"], f"{r:.4f}", strength])
            story += [tbl(rows, [5*cm, 5*cm, 3.5*cm, 3*cm]), Spacer(1, 0.5*cm)]

        # Forecasting
        fc = data.get("analyses", {}).get("forecasting", {}).get("forecasts", {})
        if fc:
            story += [hr(), Paragraph("Forecast Analysis", h1_style)]
            for col_name, fc_info in fc.items():
                pred = fc_info.get("forecast", [])
                story += [
                    Paragraph(f"Column: {col_name}", h2_style),
                    Paragraph(f"Historical data points: {fc_info.get('n_historical')}  |  "
                              f"Forecast horizon: {fc_info.get('n_forecast')} periods", body_style),
                    Paragraph("Forecasted values: " + ", ".join([f"{round(v,2)}" for v in pred if v is not None][:5]),
                              body_style),
                    Spacer(1, 0.2*cm),
                ]

        # Clustering
        cl = data.get("analyses", {}).get("clustering", {})
        summaries = cl.get("summaries", {})
        if summaries:
            story += [hr(), Paragraph(f"Cluster Analysis  (k = {cl.get('k', '?')})", h1_style)]
            for cid, cols_info in summaries.items():
                story.append(Paragraph(f"Cluster {int(cid)+1}", h2_style))
                rows = [["Column", "Mean", "Count"]]
                for col_name, stats in cols_info.items():
                    rows.append([col_name, f"{stats.get('mean') or 0:.4f}", str(stats.get("count",""))])
                story += [tbl(rows, [7*cm, 5*cm, 4*cm]), Spacer(1, 0.3*cm)]

        doc.build(story)
        buf.seek(0)
        return send_file(buf, mimetype="application/pdf",
                         as_attachment=True, download_name=f"manava_report_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf")

    except Exception:
        return jsonify({"error": traceback.format_exc()}), 500


if __name__ == "__main__":
    print("\n" + "="*52)
    print("  MANAVA Analytics Server")
    print("  http://127.0.0.1:5050")
    print("="*52 + "\n")
    app.run(host="127.0.0.1", port=5050, debug=False)
